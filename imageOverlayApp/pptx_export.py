# -*- coding: utf-8 -*-
"""Export overlay items + background image to an editable .pptx.

The background image fills one slide; every overlay becomes a real PowerPoint
text box (movable / editable) positioned to match the editor.
"""
from typing import List, Dict
from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

EMU_PER_IN = 914400
DPI = 96.0


def _rgb(hexstr: str) -> RGBColor:
    s = (hexstr or "#0f2742").lstrip("#")
    if len(s) == 3:
        s = "".join(c * 2 for c in s)
    try:
        return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except Exception:
        return RGBColor(0x0f, 0x27, 0x42)


def build_pptx(image_path: str, items: List[Dict], out_path: str) -> str:
    img = Image.open(image_path)
    Wpx, Hpx = img.size
    slide_w = Emu(int(Wpx / DPI * EMU_PER_IN))
    slide_h = Emu(int(Hpx / DPI * EMU_PER_IN))

    prs = Presentation()
    prs.slide_width = slide_w
    prs.slide_height = slide_h
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    slide.shapes.add_picture(image_path, 0, 0, width=slide_w, height=slide_h)

    align_map = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}

    for it in items:
        x = float(it.get("x", 0)); y = float(it.get("y", 0))
        w = float(it.get("w", 8)); fs = float(it.get("fs", 1.0))
        pt = max(5.0, fs * Wpx * 0.0075)
        left = Emu(int(x / 100 * slide_w))
        top = Emu(int(y / 100 * slide_h))
        width = Emu(int(max(w, 1) / 100 * slide_w))
        mh = float(it.get("mh", 0) or 0)
        if mh > 0:
            height = Emu(int(mh / 100 * slide_h))
        else:
            height = Emu(int(pt * 1.5 / 72 * EMU_PER_IN))

        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = bool(it.get("fx"))  # manually resized boxes wrap like PowerPoint
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE if mh > 0 else MSO_ANCHOR.TOP
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(tf, m, Emu(int(0.01 * EMU_PER_IN)))

        p = tf.paragraphs[0]
        p.alignment = align_map.get(it.get("a", "l"), PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = it.get("t", "")
        f = run.font
        f.size = Pt(round(pt, 1))
        f.bold = bool(it.get("b"))
        f.italic = bool(it.get("i"))
        f.underline = bool(it.get("u"))
        f.color.rgb = _rgb(it.get("c"))
        f.name = "Hiragino Sans"

        # white background mask (default on) / optional border
        if it.get("bg", 1) != 0:
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        else:
            box.fill.background()
        if it.get("bd"):
            box.line.color.rgb = RGBColor(0x9b, 0xb0, 0xc4)
            box.line.width = Pt(0.75)
        else:
            box.line.fill.background()

    prs.save(out_path)
    return out_path
